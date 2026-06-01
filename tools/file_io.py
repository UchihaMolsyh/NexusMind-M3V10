"""
File I/O — read/write multiple formats: PDF, TXT, DOCX, PPTX, XLSX, images, audio, video, etc.
"""
import os
import json
import base64
import mimetypes
from pathlib import Path
from typing import Any, Dict, Optional
from core.tool_registry import registry, ToolParam


SUPPORTED_READ = {
    ".txt", ".md", ".csv", ".json", ".xml", ".html", ".css", ".js", ".py",
    ".c", ".cpp", ".h", ".java", ".rs", ".go", ".yaml", ".yml", ".toml",
    ".ini", ".cfg", ".log", ".bat", ".ps1", ".sh",
    ".pdf", ".docx", ".pptx", ".xlsx",
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp",
    ".mp3", ".wav", ".ogg", ".flac",
    ".mp4", ".avi", ".mkv", ".webm",
    ".dll", ".exe",
}


def read_text_file(path: Path) -> str:
    encodings = ["utf-8", "latin-1", "cp1252", "ascii"]
    for enc in encodings:
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def read_pdf(path: Path) -> str:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(path))
        text = []
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text.append(t)
        return "\n\n".join(text)
    except ImportError:
        return "[PyPDF2 not installed — pip install PyPDF2]"


def read_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs)
    except ImportError:
        return "[python-docx not installed — pip install python-docx]"


def read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
            slides_text.append("\n".join(parts))
        return "\n\n".join(slides_text)
    except ImportError:
        return "[python-pptx not installed — pip install python-pptx]"


def read_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True)
        sheets = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append("\t".join(str(c) if c is not None else "" for c in row))
            sheets.append(f"=== {ws.title} ===\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets)
    except ImportError:
        return "[openpyxl not installed — pip install openpyxl]"


def read_binary_info(path: Path) -> str:
    size = path.stat().st_size
    mime = mimetypes.guess_type(str(path))[0] or "unknown"
    return f"Binary file: {path.name}\nSize: {size:,} bytes\nMIME: {mime}"


@registry.tool(
    name="read_file",
    description="Read file contents. Supports: TXT, PDF, DOCX, PPTX, XLSX, code files, images (metadata), audio/video (metadata), DLL, EXE.",
    category="File Operations",
    parameters=[
        ToolParam("path", "string", "Absolute or relative file path"),
        ToolParam("encoding", "string", "Text encoding (default: auto-detect)", required=False, default="auto"),
        ToolParam("max_chars", "string", "Max characters to return (default: 50000)", required=False, default="50000"),
    ],
)
def read_file(path: str, encoding: str = "auto", max_chars: str = "50000"):
    p = Path(path).resolve()
    if not p.exists():
        return {"error": f"File not found: {path}"}

    from config import MAX_FILE_SIZE_MB
    size_mb = p.stat().st_size / (1024 * 1024)
    if size_mb > MAX_FILE_SIZE_MB:
        return {"error": f"File too large: {size_mb:.1f}MB (max: {MAX_FILE_SIZE_MB}MB)"}

    ext = p.suffix.lower()
    max_c = int(max_chars)

    try:
        if ext == ".pdf":
            content = read_pdf(p)
        elif ext == ".docx":
            content = read_docx(p)
        elif ext == ".pptx":
            content = read_pptx(p)
        elif ext == ".xlsx":
            content = read_xlsx(p)
        elif ext in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}:
            content = read_binary_info(p)
            try:
                from PIL import Image
                img = Image.open(str(p))
                content += f"\nDimensions: {img.size[0]}x{img.size[1]}\nMode: {img.mode}"
                img.close()
            except ImportError:
                pass
        elif ext in {".mp3", ".wav", ".ogg", ".flac", ".mp4", ".avi", ".mkv", ".webm"}:
            content = read_binary_info(p)
        elif ext in {".dll", ".exe"}:
            content = read_binary_info(p)
            content += "\n[Binary executable — content not displayed]"
        else:
            content = read_text_file(p)

        if len(content) > max_c:
            content = content[:max_c] + f"\n\n... [truncated at {max_c} characters]"

        return {
            "path": str(p),
            "extension": ext,
            "size_bytes": p.stat().st_size,
            "content": content,
        }
    except Exception as e:
        return {"error": str(e), "path": str(p)}


@registry.tool(
    name="write_file",
    description="Write content to a file. Creates parent directories if needed.",
    category="File Operations",
    parameters=[
        ToolParam("path", "string", "File path to write to"),
        ToolParam("content", "string", "Content to write"),
        ToolParam("mode", "string", "Write mode: 'w' (overwrite) or 'a' (append)", required=False, default="w"),
    ],
)
def write_file(path: str, content: str, mode: str = "w"):
    p = Path(path).resolve()
    p.parent.mkdir(parents=True, exist_ok=True)

    try:
        with open(p, mode, encoding="utf-8") as f:
            f.write(content)
        return {
            "path": str(p),
            "bytes_written": len(content.encode("utf-8")),
            "mode": "overwrite" if mode == "w" else "append",
        }
    except Exception as e:
        return {"error": str(e)}


@registry.tool(
    name="list_files",
    description="List files in a directory with optional filtering.",
    category="File Operations",
    parameters=[
        ToolParam("path", "string", "Directory path"),
        ToolParam("pattern", "string", "Glob pattern filter (e.g., '*.py')", required=False, default="*"),
        ToolParam("recursive", "string", "Recurse into subdirectories (true/false)", required=False, default="false"),
    ],
)
def list_files(path: str, pattern: str = "*", recursive: str = "false"):
    p = Path(path).resolve()
    if not p.exists():
        return {"error": f"Directory not found: {path}"}
    if not p.is_dir():
        return {"error": f"Not a directory: {path}"}

    try:
        if recursive.lower() == "true":
            files = list(p.rglob(pattern))
        else:
            files = list(p.glob(pattern))

        items = []
        for f in sorted(files)[:200]:
            items.append({
                "name": f.name,
                "path": str(f),
                "is_dir": f.is_dir(),
                "size": f.stat().st_size if f.is_file() else None,
            })

        return {"directory": str(p), "count": len(items), "files": items}
    except Exception as e:
        return {"error": str(e)}
